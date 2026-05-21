"""
scripts/parse_ir_books.py
──────────────────────────
IR북 PDF/PPTX/DOCX 텍스트 추출 + IR 담당자 contact 정보 파싱.

전략:
  1) 표지(첫 1-2 페이지) + 마지막 1-2 페이지에 집중
  2) 키워드 근처 (문의처/Contact/IR/투자자정보) 컨텍스트 추출
  3) v5 정규식 재활용 (이메일/전화/이름/직책)
  4) ir_contacts에 source='IRGO_IRBOOK' 또는 'KIND_IRBOOK' 으로 저장

실행:
  python scripts/parse_ir_books.py --source IRGO          # IRGO/_metadata.json 기반 일괄
  python scripts/parse_ir_books.py --file path/to/x.pdf   # 단일 파일 디버그
"""
import io, os, re, sys, json, sqlite3, argparse
from pathlib import Path
from datetime import datetime

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

ROOT = Path(__file__).parent.parent
DB_PATH = ROOT / "data" / "dart" / "dart_reports.db"
BOOKS_DIR = ROOT / "data" / "ir_books"

KOREAN_SURNAMES = set('김이박최정강조윤장임한오서신권황안송류전홍고문양손배백허유남심노하곽성차주우구민나진지엄채원천방공현함변염여추도소석선설마길연위표명기반왕금옥육인맹제모탁국어은편복단연봉경사부')

TITLES = ['부사장','전무','상무','이사','본부장','실장','팀장','부장','차장','과장','대리','사원','주임','매니저','선임','책임','수석','Director','Manager','VP','CFO','CEO','CTO']
TITLE_RE = re.compile(r'([가-힣]{2,4})\s*[\(\[]?\s*(' + '|'.join(re.escape(t) for t in TITLES) + r')[\)\]]?')
DEPT_RE = re.compile(r'(IR|PR|홍보|투자홍보|커뮤니케이션|투자자|재무|기획|경영지원|대외협력|마케팅)\s*(팀|실|부|본부|그룹|담당)?')

EMAIL_RE = re.compile(r"[\w\._\-+]+@[\w\._\-]+\.[a-zA-Z]{2,}")
PHONE_RE = re.compile(r'(?:\+?82[-\s]?|0)?(\d{2,3})[-\s\.\)](\d{3,4})[-\s\.](\d{4})')
MOBILE_RE = re.compile(r'(?:\+?82[-\s]?|0)?(010|011|016|017|018|019)[-\s\.\)]?(\d{3,4})[-\s\.](\d{4})')

CONTACT_KEYWORDS_RE = re.compile(
    r'(IR\s*담당|IR\s*문의|투자자\s*문의|Investor\s*Relations|Contact\s*Us|Contact|문의\s*처?|연락처|투자\s*문의|IR\s*Contact|IR\s*Department)',
    re.IGNORECASE
)

EXCLUDE_PATTERNS = [
    r'@accounts\.google', r'noreply', r'no-reply', r'mailer', r'postmaster',
    r'webmaster', r'@example', r'@sample', r'@test\.',
]

FAKE_NAMES = {'대표','책임','담당','문의','각종','관련','특별','신규','오늘'}


def get_db():
    db = sqlite3.connect(str(DB_PATH), timeout=30)
    db.row_factory = sqlite3.Row
    return db


def is_valid_email(em):
    if not em:
        return False
    em = em.lower()
    for pat in EXCLUDE_PATTERNS:
        if re.search(pat, em):
            return False
    return bool(re.match(r'^[\w\._\-+]+@[\w\._\-]+\.[a-zA-Z]{2,}$', em))


def is_valid_korean_name(name):
    if not name or len(name) < 2 or len(name) > 4:
        return False
    if name in FAKE_NAMES:
        return False
    if not re.match(r'^[가-힣]+$', name):
        return False
    return name[0] in KOREAN_SURNAMES


# ── PDF 추출 ────────────────────────────────────────────────────────────────
def extract_pdf_text(path, focus_pages=True):
    """PDF에서 텍스트 추출. focus_pages=True면 처음·마지막 2페이지만"""
    try:
        import pdfplumber
        text = ''
        with pdfplumber.open(path) as pdf:
            n_pages = len(pdf.pages)
            if focus_pages and n_pages > 4:
                pages = list(range(0, 2)) + list(range(n_pages - 2, n_pages))
            else:
                pages = list(range(min(n_pages, 30)))  # 최대 30페이지
            for i in pages:
                try:
                    t = pdf.pages[i].extract_text() or ''
                    text += f'\n[Page {i+1}]\n' + t
                except Exception:
                    pass
        return text
    except Exception as e:
        return ''


def extract_pptx_text(path):
    """PPTX에서 텍스트 추출"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = ''
        # 처음·마지막 슬라이드 위주
        n = len(prs.slides)
        slides = list(prs.slides)
        if n > 4:
            target = slides[:2] + slides[-2:]
        else:
            target = slides
        for slide in target:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            text += run.text + ' '
                        text += '\n'
        return text
    except Exception as e:
        return ''


def extract_docx_text(path):
    """DOCX에서 텍스트 추출"""
    try:
        from docx import Document
        doc = Document(path)
        text = ''
        for para in doc.paragraphs:
            text += para.text + '\n'
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + ' '
                text += '\n'
        return text
    except Exception as e:
        return ''


def extract_text(path):
    """파일 형식에 따라 텍스트 추출"""
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == '.pdf':
        return extract_pdf_text(p)
    elif suffix in ('.pptx', '.ppt'):
        return extract_pptx_text(p)
    elif suffix in ('.docx', '.doc'):
        return extract_docx_text(p)
    return ''


# ── Contact 추출 ────────────────────────────────────────────────────────────
def extract_contacts_from_text(text, corp_name=''):
    """텍스트에서 IR contact 정보 추출"""
    if not text:
        return []
    contacts = []

    # contact 블록 찾기 (키워드 근처 500자)
    blocks = []
    for m in CONTACT_KEYWORDS_RE.finditer(text):
        start = max(0, m.start() - 100)
        end = min(len(text), m.end() + 800)
        blocks.append(text[start:end])

    # 키워드 못 찾으면 마지막 1500자 (보통 마지막 페이지에 disclaimer + 연락처)
    if not blocks:
        blocks.append(text[-1500:])
        # 추가로 첫 1000자 (표지에 연락처 있을 수도)
        blocks.append(text[:1000])

    seen_emails = set()
    for block in blocks:
        # 이메일 추출
        emails = [e for e in EMAIL_RE.findall(block) if is_valid_email(e)]
        # 이름·직책
        name_titles = [(n, t) for n, t in TITLE_RE.findall(block) if is_valid_korean_name(n)]
        # 부서
        dept_m = DEPT_RE.search(block)
        dept = dept_m.group(1) if dept_m else ''
        # 휴대폰
        mob_m = MOBILE_RE.search(block)
        mobile = f"0{mob_m.group(1)}-{mob_m.group(2)}-{mob_m.group(3)}" if mob_m else ''
        # 일반 전화
        phones = PHONE_RE.findall(block)
        phone = ''
        for p in phones:
            prefix = p[0] if p[0].startswith('0') else '0' + p[0]
            cand = f"{prefix}-{p[1]}-{p[2]}"
            if cand != mobile:
                phone = cand
                break

        # 합치기 - 이메일 1개당 1 contact
        for em in emails:
            if em in seen_emails:
                continue
            seen_emails.add(em)
            name = name_titles[0][0] if name_titles else ''
            title = name_titles[0][1] if name_titles else ''
            contacts.append({
                'email': em,
                'name': name,
                'title': title,
                'dept': dept,
                'phone': phone,
                'mobile': mobile,
            })

        # 이메일 없어도 이름+전화면 등록
        if not emails and name_titles and (mobile or phone):
            contacts.append({
                'email': '',
                'name': name_titles[0][0],
                'title': name_titles[0][1],
                'dept': dept,
                'phone': phone,
                'mobile': mobile,
            })

    # 중복 제거
    unique = []
    seen = set()
    for c in contacts:
        key = (c['email'], c['name'], c['phone'] or c['mobile'])
        if key not in seen:
            seen.add(key)
            unique.append(c)
    return unique


# ── DB 저장 ──────────────────────────────────────────────────────────────────
def save_to_db(db, corp_code, corp_name, contact, source, source_url, file_path):
    """ir_contacts에 저장"""
    if not (contact.get('email') or contact.get('name')):
        return False
    now = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

    # 신뢰도 결정
    if contact.get('email'):
        if contact.get('name') and (contact.get('phone') or contact.get('mobile')):
            conf = 'A_complete'
        else:
            conf = 'A_email_only'
    else:
        conf = 'B_phone_only'

    email = contact.get('email') or None
    unique_email = email if email else f"_irbook:{corp_code}:{contact.get('mobile') or contact.get('phone')}:{contact.get('name')}"

    try:
        if email:
            db.execute("""
                INSERT INTO ir_contacts
                    (corp_code, corp_name, ir_email, ir_phone, ir_mobile,
                     ir_name, ir_dept, ir_title, source, source_url,
                     confidence, mx_verified, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, 1, ?, ?)
                ON CONFLICT(corp_code, ir_email) DO UPDATE SET
                    ir_phone = COALESCE(NULLIF(excluded.ir_phone,''), ir_contacts.ir_phone),
                    ir_mobile = COALESCE(NULLIF(excluded.ir_mobile,''), ir_contacts.ir_mobile),
                    ir_name = COALESCE(NULLIF(excluded.ir_name,''), ir_contacts.ir_name),
                    ir_dept = COALESCE(NULLIF(excluded.ir_dept,''), ir_contacts.ir_dept),
                    ir_title = COALESCE(NULLIF(excluded.ir_title,''), ir_contacts.ir_title),
                    confidence = excluded.confidence,
                    cross_verified = ir_contacts.cross_verified + 1,
                    updated_at = excluded.updated_at
            """, [corp_code, corp_name, email, contact.get('phone',''), contact.get('mobile',''),
                  contact.get('name',''), contact.get('dept',''), contact.get('title',''),
                  source, source_url, conf, now, now])
        else:
            # 이메일 없는 경우
            db.execute("""
                INSERT OR IGNORE INTO ir_contacts
                    (corp_code, corp_name, ir_phone, ir_mobile, ir_name, ir_dept, ir_title,
                     source, source_url, confidence, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, [corp_code, corp_name, contact.get('phone',''), contact.get('mobile',''),
                  contact.get('name',''), contact.get('dept',''), contact.get('title',''),
                  source, source_url, conf, now, now])
        db.commit()
        return True
    except Exception as e:
        return False


# ── 메인 ─────────────────────────────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--source', choices=['IRGO', 'KIND'], default='IRGO')
    ap.add_argument('--file', type=str, help='단일 파일 디버그')
    ap.add_argument('--limit', type=int, default=0, help='처리 최대 (0=전부)')
    args = ap.parse_args()

    db = get_db()

    if args.file:
        # 단일 파일 디버그
        text = extract_text(args.file)
        print(f'추출 텍스트 길이: {len(text)}')
        if text:
            contacts = extract_contacts_from_text(text)
            print(f'추출 contact: {len(contacts)}건')
            for c in contacts:
                print(f'  {c}')
        return

    # IRGO 메타데이터 기반 일괄
    meta_file = BOOKS_DIR / args.source / '_metadata.json'
    if not meta_file.exists():
        print(f'❌ {meta_file} 없음. 먼저 collect_*.py 로 다운로드 필요.')
        return

    meta = json.loads(meta_file.read_text(encoding='utf-8'))
    downloaded = meta.get('downloaded', {})
    matched = meta.get('matched_companies', {})

    print('=' * 70)
    print(f'  IR북 파싱 시작 — source: {args.source}')
    print('=' * 70)
    print(f'  대상 파일: {sum(len(v) for v in downloaded.values())}개')
    print(f'  대상 회사: {len(matched)}개\n')

    stats = {'parsed': 0, 'with_email': 0, 'with_phone': 0, 'saved': 0,
             'no_text': 0, 'no_contact': 0, 'errors': 0}
    source_label = f'{args.source}_IRBOOK'

    items = []
    for corp_code, files in downloaded.items():
        for f in files:
            items.append((corp_code, f))

    if args.limit > 0:
        items = items[:args.limit]

    for i, (corp_code, finfo) in enumerate(items, 1):
        corp_name = matched.get(corp_code, '')
        pdf_path = ROOT / finfo['pdf_path']
        if not pdf_path.exists():
            stats['errors'] += 1
            continue

        try:
            text = extract_text(pdf_path)
            stats['parsed'] += 1

            if not text or len(text) < 50:
                stats['no_text'] += 1
                continue

            contacts = extract_contacts_from_text(text, corp_name)
            if not contacts:
                stats['no_contact'] += 1
                continue

            for c in contacts:
                if c.get('email'):
                    stats['with_email'] += 1
                elif c.get('phone') or c.get('mobile'):
                    stats['with_phone'] += 1
                if save_to_db(db, corp_code, corp_name, c,
                              source_label, f'{args.source.lower()}:{finfo.get("irgo_id","")}',
                              str(pdf_path)):
                    stats['saved'] += 1

            if i % 10 == 0:
                print(f'  [{i:4d}/{len(items)}] 파싱 중 — '
                      f'이메일 {stats["with_email"]} | 전화 {stats["with_phone"]} | 저장 {stats["saved"]}')
        except Exception as e:
            stats['errors'] += 1
            if stats['errors'] <= 3:
                print(f'  ERR {pdf_path.name}: {str(e)[:80]}')

    print('\n' + '=' * 70)
    print(f'  파싱 완료')
    print('=' * 70)
    print(f'  파싱 시도        : {stats["parsed"]}')
    print(f'  ⚠ 텍스트 없음    : {stats["no_text"]} (이미지 PDF 등)')
    print(f'  ⚠ contact 없음   : {stats["no_contact"]}')
    print(f'  ✅ 이메일 추출   : {stats["with_email"]}')
    print(f'  ✅ 전화만 추출   : {stats["with_phone"]}')
    print(f'  💾 DB 저장       : {stats["saved"]}')
    print(f'  ❌ 에러          : {stats["errors"]}')

    # 누적
    n = db.execute(f"SELECT COUNT(*) FROM ir_contacts WHERE source='{source_label}'").fetchone()[0]
    n_corp = db.execute(f"SELECT COUNT(DISTINCT corp_code) FROM ir_contacts WHERE source='{source_label}' AND corp_code != 'UNKNOWN'").fetchone()[0]
    print(f'\n  [{source_label}] 누적: {n}건 / {n_corp}개 회사')


if __name__ == '__main__':
    main()
